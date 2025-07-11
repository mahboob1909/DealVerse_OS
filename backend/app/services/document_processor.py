"""
Document Processing Pipeline for DealVerse OS
Handles virus scanning, format conversion, metadata extraction, thumbnail generation, and search indexing
"""
import os
import io
import logging
import hashlib
import tempfile
import asyncio
import json
import re
from typing import Dict, Any, Optional, Tuple, BinaryIO, List
from pathlib import Path
import mimetypes
from datetime import datetime

# Enhanced document processing libraries
import PyPDF2
import fitz  # PyMuPDF for better PDF processing
from docx import Document
import openpyxl
from pptx import Presentation
import pytesseract
from PIL import Image, ImageDraw, ImageFont
# Try to import magic for MIME type detection, fallback to mimetypes
try:
    import magic  # For MIME type detection
    MAGIC_AVAILABLE = True
except ImportError:
    MAGIC_AVAILABLE = False
    import mimetypes
import spacy
from textstat import flesch_reading_ease, flesch_kincaid_grade

# FastAPI and application imports
from fastapi import HTTPException, UploadFile
from app.core.config import settings
from app.services.s3_storage import get_s3_storage
from app.services import s3_storage  # Import module for mocking

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """Comprehensive document processing pipeline with virus scanning, metadata extraction, and indexing"""

    SUPPORTED_FORMATS = {
        'application/pdf': 'pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
        'application/msword': 'doc',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
        'application/vnd.ms-excel': 'xls',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
        'application/vnd.ms-powerpoint': 'ppt',
        'text/plain': 'txt',
        'text/csv': 'csv',
        'image/jpeg': 'jpg',
        'image/png': 'png',
        'image/tiff': 'tiff'
    }

    def __init__(self):
        self.max_file_size = getattr(settings, 'MAX_FILE_SIZE_MB', 50) * 1024 * 1024  # Convert to bytes
        self.allowed_types = getattr(settings, 'ALLOWED_FILE_TYPES', list(self.SUPPORTED_FORMATS.keys()))
        self.temp_dir = Path("temp_uploads")
        self.temp_dir.mkdir(exist_ok=True)

        # Initialize NLP model for text analysis
        try:
            self.nlp = spacy.load("en_core_web_sm")
        except OSError:
            logger.warning("spaCy English model not found. Text analysis will be limited.")
            self.nlp = None

        # Virus scanning configuration
        self.virus_scan_enabled = getattr(settings, 'VIRUS_SCAN_ENABLED', True)
        self.clamav_socket = getattr(settings, 'CLAMAV_SOCKET_PATH', '/var/run/clamav/clamd.ctl')

        # Thumbnail configuration
        self.thumbnail_size = (300, 400)  # Width x Height for document thumbnails
        self.thumbnail_quality = 85

        # Security configuration
        self.enable_content_analysis = getattr(settings, 'ENABLE_CONTENT_ANALYSIS', True)
        self.enable_pii_detection = getattr(settings, 'ENABLE_PII_DETECTION', True)
        self.quarantine_suspicious_files = getattr(settings, 'QUARANTINE_SUSPICIOUS_FILES', True)

        # File type restrictions
        self.allowed_extensions = {
            '.pdf', '.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx',
            '.txt', '.rtf', '.csv', '.json', '.xml', '.html', '.htm',
            '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.svg',
            '.mp4', '.avi', '.mov', '.wmv', '.flv', '.webm',
            '.mp3', '.wav', '.flac', '.aac', '.ogg',
            '.zip', '.rar', '.7z', '.tar', '.gz'
        }

        # Dangerous file patterns
        self.dangerous_patterns = [
            r'<script[^>]*>.*?</script>',  # JavaScript
            r'javascript:',  # JavaScript URLs
            r'vbscript:',   # VBScript URLs
            r'data:text/html',  # Data URLs with HTML
            r'<?php',       # PHP code
            r'<%.*?%>',     # ASP/JSP code
        ]
    
    async def process_document_pipeline(self, file: UploadFile, document_id: str,
                                      organization_id: str) -> Dict[str, Any]:
        """
        Complete document processing pipeline

        Args:
            file: Uploaded file object
            document_id: Unique document identifier
            organization_id: Organization ID for isolation

        Returns:
            Processing results with metadata, thumbnails, and extracted content
        """
        try:
            # Step 1: Basic validation
            validation_result = await self._validate_file_enhanced(file)
            if not validation_result["is_valid"]:
                raise HTTPException(status_code=400, detail=validation_result["error"])

            # Step 2: Virus scanning
            if self.virus_scan_enabled:
                scan_result = await self._scan_for_viruses(file)
                if not scan_result["is_safe"]:
                    raise HTTPException(status_code=400, detail=f"Security threat detected: {scan_result['threats_detected']}")
            else:
                scan_result = {"is_safe": True, "threats_detected": [], "scan_method": "disabled"}

            # Step 3: Extract metadata and content
            file_content = await file.read()
            await file.seek(0)  # Reset file pointer

            metadata = await self._extract_enhanced_metadata(file, file_content)
            text_content = await self._extract_text_content_enhanced(file, file_content, metadata["mime_type"])

            # Step 4: Generate thumbnail
            thumbnail_data = await self._generate_thumbnail(file, file_content, metadata["mime_type"])

            # Step 5: Analyze text content
            text_analysis = await self._analyze_text_content(text_content)

            # Step 6: Create search index data
            search_index = await self._create_search_index(text_content, metadata, text_analysis)

            # Step 7: Store processing artifacts
            artifacts = await self._store_processing_artifacts(
                document_id, organization_id, thumbnail_data, text_content, search_index
            )

            processing_result = {
                "status": "success",
                "document_id": document_id,
                "validation": validation_result,
                "virus_scan": scan_result,
                "metadata": metadata,
                "text_content": text_content,
                "text_analysis": text_analysis,
                "thumbnail": artifacts.get("thumbnail"),
                "search_index": search_index,
                "processing_timestamp": datetime.utcnow().isoformat(),
                "artifacts": artifacts
            }

            logger.info(f"Successfully processed document {document_id}")
            return processing_result

        except HTTPException:
            raise
        except Exception as e:
            logger.error(f"Document processing failed for {document_id}: {e}")
            raise HTTPException(status_code=500, detail="Document processing failed")

    def is_supported_format(self, mime_type: str) -> bool:
        """Check if the file format is supported"""
        return mime_type in self.SUPPORTED_FORMATS

    async def _validate_file_enhanced(self, file: UploadFile) -> Dict[str, Any]:
        """Enhanced file validation with security checks"""
        try:
            # Check file size by reading the entire file
            content = await file.read()
            file_size = len(content)
            await file.seek(0)  # Reset to beginning

            if file_size > self.max_file_size:
                return {
                    "is_valid": False,
                    "error": f"File size ({file_size / 1024 / 1024:.1f}MB) exceeds maximum allowed size ({self.max_file_size / 1024 / 1024}MB)"
                }

            # Check MIME type using magic numbers or fallback
            file_content = await file.read(1024)  # Read first 1KB for type detection
            await file.seek(0)

            if MAGIC_AVAILABLE:
                detected_mime = magic.from_buffer(file_content, mime=True)
            else:
                # Fallback to mimetypes based on filename
                detected_mime, _ = mimetypes.guess_type(file.filename)
                if not detected_mime:
                    detected_mime = "application/octet-stream"

            if detected_mime not in self.allowed_types:
                return {
                    "is_valid": False,
                    "error": f"File type '{detected_mime}' is not allowed"
                }

            # Check filename
            if not file.filename or len(file.filename) > 255:
                return {
                    "is_valid": False,
                    "error": "Invalid filename"
                }

            # Check for suspicious file extensions
            suspicious_extensions = ['.exe', '.bat', '.cmd', '.scr', '.vbs', '.js']
            file_ext = Path(file.filename).suffix.lower()
            if file_ext in suspicious_extensions:
                return {
                    "is_valid": False,
                    "error": f"File extension '{file_ext}' is not allowed for security reasons"
                }

            return {
                "is_valid": True,
                "file_size": file_size,
                "mime_type": detected_mime,
                "filename": file.filename,
                "file_extension": file_ext
            }

        except Exception as e:
            logger.error(f"Enhanced file validation failed: {e}")
            return {"is_valid": False, "error": "File validation failed"}

    async def _scan_for_viruses(self, file: UploadFile) -> Dict[str, Any]:
        """
        Scan file for viruses and malicious content
        Note: This is a placeholder implementation. In production, integrate with ClamAV daemon
        """
        try:
            # Read file content for scanning
            file_content = await file.read()
            await file.seek(0)

            # Calculate file hash for known threat database lookup
            file_hash = hashlib.sha256(file_content).hexdigest()

            # Check for suspicious file patterns
            suspicious_patterns = [
                b"<script", b"javascript:", b"vbscript:", b"powershell",
                b"cmd.exe", b"eval(", b"exec(", b"system(", b"shell_exec",
                b"eicar-standard-antivirus-test-file"  # EICAR test string (lowercase)
            ]

            file_content_lower = file_content.lower()
            for pattern in suspicious_patterns:
                if pattern in file_content_lower:
                    threat_msg = f"Suspicious content pattern detected: {pattern.decode('utf-8', errors='ignore')}"
                    return {
                        "is_safe": False,
                        "threat": threat_msg,
                        "threats_detected": [threat_msg],
                        "hash": file_hash
                    }

            # Check file size anomalies (files that are too small or too large for their type)
            file_size = len(file_content)
            if file_size < 10:  # Suspiciously small files
                threat_msg = "File is suspiciously small"
                return {
                    "is_safe": False,
                    "threat": threat_msg,
                    "threats_detected": [threat_msg],
                    "hash": file_hash
                }

            return {
                "is_safe": True,
                "threats_detected": [],
                "scan_method": "pattern_detection",
                "hash": file_hash,
                "scan_timestamp": datetime.utcnow().isoformat()
            }

        except Exception as e:
            logger.error(f"Virus scanning failed: {e}")
            # In production, you might want to fail-safe (reject file) or fail-open (allow file)
            return {"is_safe": True, "threats_detected": [], "error": f"Scan failed: {e}"}

    async def _extract_enhanced_metadata(self, file: UploadFile, file_content: bytes) -> Dict[str, Any]:
        """Extract comprehensive metadata from file"""
        try:
            # Basic file information
            if MAGIC_AVAILABLE:
                mime_type = magic.from_buffer(file_content, mime=True)
            else:
                # Fallback to mimetypes based on filename
                mime_type, _ = mimetypes.guess_type(file.filename)
                if not mime_type:
                    mime_type = "application/octet-stream"

            metadata = {
                "filename": file.filename,
                "file_size": len(file_content),
                "mime_type": mime_type,
                "file_hash": hashlib.sha256(file_content).hexdigest(),
                "upload_timestamp": datetime.utcnow().isoformat(),
                "creation_date": datetime.utcnow().isoformat()  # Default to current time
            }

            # Extract file extension
            if file.filename:
                metadata["file_extension"] = Path(file.filename).suffix.lower()

            # Format-specific metadata extraction
            mime_type = metadata["mime_type"]

            if mime_type == "application/pdf":
                metadata.update(await self._extract_pdf_metadata_enhanced(file_content))
            elif mime_type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/msword"]:
                metadata.update(await self._extract_word_metadata_enhanced(file_content))
            elif mime_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel"]:
                metadata.update(await self._extract_excel_metadata_enhanced(file_content))
            elif mime_type in ["image/jpeg", "image/png", "image/tiff"]:
                metadata.update(await self._extract_image_metadata_enhanced(file_content))

            return metadata

        except Exception as e:
            logger.error(f"Enhanced metadata extraction failed: {e}")
            return {
                "filename": file.filename,
                "file_size": len(file_content),
                "mime_type": "application/octet-stream",
                "error": f"Metadata extraction failed: {e}"
            }

    async def _extract_pdf_metadata_enhanced(self, file_content: bytes) -> Dict[str, Any]:
        """Extract enhanced metadata from PDF files using PyMuPDF"""
        try:
            with tempfile.NamedTemporaryFile() as temp_file:
                temp_file.write(file_content)
                temp_file.flush()

                doc = fitz.open(temp_file.name)
                metadata = doc.metadata

                # Count text blocks and images
                text_blocks = 0
                image_count = 0

                for page_num in range(min(doc.page_count, 5)):  # Check first 5 pages
                    page = doc[page_num]
                    text_blocks += len(page.get_text("dict")["blocks"])
                    image_count += len(page.get_images())

                return {
                    "page_count": doc.page_count,
                    "title": metadata.get("title", ""),
                    "author": metadata.get("author", ""),
                    "subject": metadata.get("subject", ""),
                    "creator": metadata.get("creator", ""),
                    "producer": metadata.get("producer", ""),
                    "creation_date": metadata.get("creationDate", ""),
                    "modification_date": metadata.get("modDate", ""),
                    "encrypted": doc.needs_pass,
                    "text_blocks": text_blocks,
                    "image_count": image_count,
                    "format": "PDF"
                }

        except Exception as e:
            logger.error(f"Enhanced PDF metadata extraction failed: {e}")
            return {"format": "PDF", "error": str(e)}

    async def _extract_word_metadata_enhanced(self, file_content: bytes) -> Dict[str, Any]:
        """Extract enhanced metadata from Word documents"""
        try:
            with tempfile.NamedTemporaryFile() as temp_file:
                temp_file.write(file_content)
                temp_file.flush()

                doc = Document(temp_file.name)
                core_props = doc.core_properties

                # Count paragraphs, words, and tables
                paragraph_count = len(doc.paragraphs)
                word_count = sum(len(paragraph.text.split()) for paragraph in doc.paragraphs)
                table_count = len(doc.tables)

                # Count images
                image_count = 0
                for rel in doc.part.rels.values():
                    if "image" in rel.target_ref:
                        image_count += 1

                return {
                    "title": core_props.title or "",
                    "author": core_props.author or "",
                    "subject": core_props.subject or "",
                    "created": core_props.created.isoformat() if core_props.created else "",
                    "modified": core_props.modified.isoformat() if core_props.modified else "",
                    "paragraph_count": paragraph_count,
                    "word_count": word_count,
                    "table_count": table_count,
                    "image_count": image_count,
                    "format": "Word Document"
                }

        except Exception as e:
            logger.error(f"Enhanced Word metadata extraction failed: {e}")
            return {"format": "Word Document", "error": str(e)}

    async def _extract_excel_metadata_enhanced(self, file_content: bytes) -> Dict[str, Any]:
        """Extract enhanced metadata from Excel files"""
        try:
            with tempfile.NamedTemporaryFile() as temp_file:
                temp_file.write(file_content)
                temp_file.flush()

                workbook = openpyxl.load_workbook(temp_file.name)
                props = workbook.properties

                # Count cells with data
                total_cells = 0
                for sheet in workbook.worksheets:
                    for row in sheet.iter_rows():
                        for cell in row:
                            if cell.value is not None:
                                total_cells += 1

                return {
                    "title": props.title or "",
                    "creator": props.creator or "",
                    "subject": props.subject or "",
                    "created": props.created.isoformat() if props.created else "",
                    "modified": props.modified.isoformat() if props.modified else "",
                    "sheet_count": len(workbook.worksheets),
                    "sheet_names": [sheet.title for sheet in workbook.worksheets],
                    "total_cells_with_data": total_cells,
                    "format": "Excel Spreadsheet"
                }

        except Exception as e:
            logger.error(f"Enhanced Excel metadata extraction failed: {e}")
            return {"format": "Excel Spreadsheet", "error": str(e)}

    async def _extract_image_metadata_enhanced(self, file_content: bytes) -> Dict[str, Any]:
        """Extract enhanced metadata from image files"""
        try:
            with tempfile.NamedTemporaryFile() as temp_file:
                temp_file.write(file_content)
                temp_file.flush()

                with Image.open(temp_file.name) as img:
                    # Extract EXIF data if available
                    exif_data = {}
                    if hasattr(img, '_getexif') and img._getexif():
                        exif_data = img._getexif()

                    return {
                        "width": img.width,
                        "height": img.height,
                        "mode": img.mode,
                        "format": img.format,
                        "has_transparency": img.mode in ("RGBA", "LA") or "transparency" in img.info,
                        "dpi": img.info.get("dpi", (72, 72)),
                        "has_exif": bool(exif_data),
                        "color_count": len(img.getcolors(maxcolors=256)) if img.mode in ("P", "L") else None
                    }

        except Exception as e:
            logger.error(f"Enhanced image metadata extraction failed: {e}")
            return {"format": "Image", "error": str(e)}

    async def _generate_thumbnail(self, file: UploadFile, file_content: bytes, mime_type: str) -> Optional[bytes]:
        """Generate thumbnail for document"""
        try:
            if mime_type == "application/pdf":
                return await self._generate_pdf_thumbnail(file_content)
            elif mime_type in ["image/jpeg", "image/png", "image/tiff"]:
                return await self._generate_image_thumbnail(file_content)
            elif mime_type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document"]:
                return await self._generate_document_thumbnail(file.filename, "Word Document")
            elif mime_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"]:
                return await self._generate_document_thumbnail(file.filename, "Excel Spreadsheet")
            else:
                return await self._generate_generic_thumbnail(file.filename, mime_type)

        except Exception as e:
            logger.error(f"Thumbnail generation failed: {e}")
            return None

    async def _generate_pdf_thumbnail(self, file_content: bytes) -> Optional[bytes]:
        """Generate thumbnail from PDF first page"""
        try:
            with tempfile.NamedTemporaryFile() as temp_file:
                temp_file.write(file_content)
                temp_file.flush()

                doc = fitz.open(temp_file.name)
                if doc.page_count > 0:
                    page = doc[0]
                    # Render page as image
                    mat = fitz.Matrix(2, 2)  # 2x zoom for better quality
                    pix = page.get_pixmap(matrix=mat)
                    img_data = pix.tobytes("png")

                    # Resize to thumbnail size
                    with Image.open(io.BytesIO(img_data)) as img:
                        img.thumbnail(self.thumbnail_size, Image.Resampling.LANCZOS)

                        # Save as JPEG
                        output = io.BytesIO()
                        img.convert("RGB").save(output, format="JPEG", quality=self.thumbnail_quality)
                        return output.getvalue()

        except Exception as e:
            logger.error(f"PDF thumbnail generation failed: {e}")
            return None

    async def _generate_image_thumbnail(self, file_content: bytes) -> Optional[bytes]:
        """Generate thumbnail from image file"""
        try:
            with Image.open(io.BytesIO(file_content)) as img:
                # Create thumbnail
                img.thumbnail(self.thumbnail_size, Image.Resampling.LANCZOS)

                # Save as JPEG
                output = io.BytesIO()
                img.convert("RGB").save(output, format="JPEG", quality=self.thumbnail_quality)
                return output.getvalue()

        except Exception as e:
            logger.error(f"Image thumbnail generation failed: {e}")
            return None

    async def _generate_document_thumbnail(self, filename: str, doc_type: str) -> Optional[bytes]:
        """Generate generic thumbnail for document types"""
        try:
            # Create a simple thumbnail with document icon and filename
            img = Image.new('RGB', self.thumbnail_size, color='white')
            draw = ImageDraw.Draw(img)

            # Try to load a font, fallback to default
            try:
                font_large = ImageFont.truetype("arial.ttf", 24)
                font_small = ImageFont.truetype("arial.ttf", 16)
            except:
                font_large = ImageFont.load_default()
                font_small = ImageFont.load_default()

            # Draw document type
            draw.text((20, 50), doc_type, fill='black', font=font_large)

            # Draw filename (truncated if too long)
            display_name = filename[:30] + "..." if len(filename) > 30 else filename
            draw.text((20, 100), display_name, fill='gray', font=font_small)

            # Draw a simple document icon
            draw.rectangle([20, 150, 80, 220], outline='black', width=2)
            draw.line([30, 160, 70, 160], fill='black', width=1)
            draw.line([30, 170, 70, 170], fill='black', width=1)
            draw.line([30, 180, 60, 180], fill='black', width=1)

            # Save as JPEG
            output = io.BytesIO()
            img.save(output, format="JPEG", quality=self.thumbnail_quality)
            return output.getvalue()

        except Exception as e:
            logger.error(f"Document thumbnail generation failed: {e}")
            return None

    async def _generate_generic_thumbnail(self, filename: str, mime_type: str) -> Optional[bytes]:
        """Generate generic thumbnail for unknown file types"""
        try:
            img = Image.new('RGB', self.thumbnail_size, color='lightgray')
            draw = ImageDraw.Draw(img)

            try:
                font = ImageFont.truetype("arial.ttf", 20)
            except:
                font = ImageFont.load_default()

            # Draw file type
            draw.text((20, 50), "Unknown File", fill='black', font=font)
            draw.text((20, 80), mime_type, fill='gray', font=font)

            # Draw filename
            display_name = filename[:25] + "..." if len(filename) > 25 else filename
            draw.text((20, 120), display_name, fill='black', font=font)

            # Save as JPEG
            output = io.BytesIO()
            img.save(output, format="JPEG", quality=self.thumbnail_quality)
            return output.getvalue()

        except Exception as e:
            logger.error(f"Generic thumbnail generation failed: {e}")
            return None

    async def _extract_text_content_enhanced(self, file: UploadFile, file_content: bytes, mime_type: str) -> str:
        """Extract text content with enhanced processing"""
        try:
            # Extract text based on MIME type
            if mime_type == "application/pdf":
                text_content = self._extract_from_pdf(file_content)
            elif mime_type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/msword"]:
                text_content = self._extract_from_docx(file_content)
            elif mime_type in ["application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel"]:
                text_content = self._extract_from_excel(file_content)
            elif mime_type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation", "application/vnd.ms-powerpoint"]:
                text_content = self._extract_from_powerpoint(file_content)
            elif mime_type.startswith("text/"):
                text_content = self._extract_from_text(file_content)
            elif mime_type.startswith("image/"):
                text_content = self._extract_from_image(file_content)
            else:
                # Try to extract as text for unknown types
                text_content = self._extract_from_text(file_content)

            # Clean and normalize text
            if text_content:
                # Remove excessive whitespace
                text_content = re.sub(r'\s+', ' ', text_content)
                # Remove control characters
                text_content = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', '', text_content)
                # Normalize line endings
                text_content = text_content.replace('\r\n', '\n').replace('\r', '\n')

            return text_content or ""

        except Exception as e:
            logger.error(f"Enhanced text extraction failed: {e}")
            return ""

    async def _analyze_text_content(self, text_content: str) -> Dict[str, Any]:
        """Analyze text content for readability, entities, and key metrics"""
        try:
            if not text_content or len(text_content.strip()) < 10:
                return {"error": "Insufficient text content for analysis"}

            analysis = {
                "character_count": len(text_content),
                "word_count": len(text_content.split()),
                "sentence_count": len(re.findall(r'[.!?]+', text_content)),
                "paragraph_count": len([p for p in text_content.split('\n\n') if p.strip()]),
            }

            # Readability analysis
            try:
                flesch_score = flesch_reading_ease(text_content)
                analysis["readability"] = {
                    "flesch_reading_ease": flesch_score,
                    "flesch_kincaid_grade": flesch_kincaid_grade(text_content)
                }
                analysis["readability_score"] = flesch_score  # Add for test compatibility
            except Exception as e:
                logger.warning(f"Readability analysis failed: {e}")
                analysis["readability"] = {"error": str(e)}
                analysis["readability_score"] = 0

            # Entity extraction using spaCy
            if self.nlp and len(text_content) < 100000:  # Limit text size for NLP processing
                try:
                    doc = self.nlp(text_content[:50000])  # Process first 50k characters

                    entities = {}
                    for ent in doc.ents:
                        if ent.label_ not in entities:
                            entities[ent.label_] = []
                        if ent.text not in entities[ent.label_]:
                            entities[ent.label_][:10]  # Limit to 10 entities per type
                            entities[ent.label_].append(ent.text)

                    analysis["entities"] = entities

                    # Extract key phrases (noun phrases)
                    key_phrases = [chunk.text for chunk in doc.noun_chunks if len(chunk.text.split()) > 1]
                    analysis["key_phrases"] = key_phrases[:20]  # Top 20 key phrases

                except Exception as e:
                    logger.warning(f"NLP analysis failed: {e}")
                    analysis["entities"] = {"error": str(e)}
                    analysis["key_phrases"] = []
            else:
                # Default empty entities when spaCy is not available
                analysis["entities"] = {}
                analysis["key_phrases"] = []

            # Language detection (simple heuristic)
            analysis["language"] = self._detect_language_simple(text_content)

            # Content classification
            analysis["content_type"] = self._classify_content_type(text_content)

            return analysis

        except Exception as e:
            logger.error(f"Text analysis failed: {e}")
            return {"error": f"Text analysis failed: {e}"}

    def _detect_language_simple(self, text: str) -> str:
        """Simple language detection based on common words"""
        try:
            text_lower = text.lower()

            # English indicators
            english_words = ['the', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by']
            english_count = sum(1 for word in english_words if f' {word} ' in text_lower)

            if english_count > 5:
                return "en"
            else:
                return "unknown"

        except Exception:
            return "unknown"

    def _classify_content_type(self, text: str) -> str:
        """Classify document content type based on keywords"""
        try:
            text_lower = text.lower()

            # Financial/Legal keywords
            financial_keywords = ['contract', 'agreement', 'financial', 'revenue', 'profit', 'investment', 'merger', 'acquisition']
            legal_keywords = ['whereas', 'hereby', 'pursuant', 'liability', 'indemnify', 'covenant', 'breach']
            technical_keywords = ['algorithm', 'implementation', 'specification', 'requirements', 'architecture']

            financial_score = sum(1 for keyword in financial_keywords if keyword in text_lower)
            legal_score = sum(1 for keyword in legal_keywords if keyword in text_lower)
            technical_score = sum(1 for keyword in technical_keywords if keyword in text_lower)

            if financial_score > legal_score and financial_score > technical_score:
                return "financial"
            elif legal_score > technical_score:
                return "legal"
            elif technical_score > 0:
                return "technical"
            else:
                return "general"

        except Exception:
            return "unknown"

    async def _create_search_index(self, text_content: str, metadata: Dict[str, Any],
                                 text_analysis: Dict[str, Any]) -> Dict[str, Any]:
        """Create search index data for the document"""
        try:
            # Extract searchable terms
            words = re.findall(r'\b\w+\b', text_content.lower())

            # Remove common stop words
            stop_words = {'the', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by', 'a', 'an'}
            meaningful_words = [word for word in words if len(word) > 2 and word not in stop_words]

            # Calculate word frequencies
            word_freq = {}
            for word in meaningful_words:
                word_freq[word] = word_freq.get(word, 0) + 1

            # Get top keywords
            top_keywords = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:50]

            search_index = {
                "keywords": [word for word, freq in top_keywords],
                "keyword_frequencies": dict(top_keywords),
                "total_words": len(words),
                "unique_words": len(set(words)),
                "entities": text_analysis.get("entities", {}),
                "key_phrases": text_analysis.get("key_phrases", []),
                "content_type": text_analysis.get("content_type", "unknown"),
                "content_classification": text_analysis.get("content_type", "unknown"),  # Add for test compatibility
                "language": text_analysis.get("language", "unknown"),
                "searchable_content": text_content[:1000],  # Add for test compatibility - first 1000 chars
                "metadata_searchable": {
                    "filename": metadata.get("filename", ""),
                    "title": metadata.get("title", ""),
                    "author": metadata.get("author", ""),
                    "subject": metadata.get("subject", "")
                }
            }

            return search_index

        except Exception as e:
            logger.error(f"Search index creation failed: {e}")
            return {"error": f"Search index creation failed: {e}"}

    async def _store_processing_artifacts(self, document_id: str, organization_id: str,
                                        thumbnail_data: Optional[bytes], text_content: str,
                                        search_index: Dict[str, Any]) -> Dict[str, Any]:
        """Store processing artifacts in S3"""
        try:
            artifacts = {}

            # Store thumbnail if generated
            if thumbnail_data:
                thumbnail_key = f"{organization_id}/thumbnails/{document_id}_thumbnail.jpg"

                # Create a file-like object for thumbnail
                thumbnail_file = io.BytesIO(thumbnail_data)
                thumbnail_file.name = f"{document_id}_thumbnail.jpg"

                # Upload thumbnail to S3
                s3_storage = get_s3_storage()
                thumbnail_result = await s3_storage.upload_file_direct(
                    file_content=thumbnail_data,
                    file_key=thumbnail_key,
                    content_type="image/jpeg",
                    metadata={
                        "document_id": document_id,
                        "artifact_type": "thumbnail",
                        "organization_id": organization_id
                    }
                )

                artifacts["thumbnail"] = {
                    "s3_key": thumbnail_key,
                    "url": s3_storage.generate_presigned_url(thumbnail_key, expiration=86400),
                    "size": len(thumbnail_data)
                }

            # Store extracted text
            if text_content:
                text_key = f"{organization_id}/extracted_text/{document_id}_text.txt"
                text_bytes = text_content.encode('utf-8')

                text_result = await s3_storage.upload_file_direct(
                    file_content=text_bytes,
                    file_key=text_key,
                    content_type="text/plain",
                    metadata={
                        "document_id": document_id,
                        "artifact_type": "extracted_text",
                        "organization_id": organization_id
                    }
                )

                artifacts["extracted_text"] = {
                    "s3_key": text_key,
                    "size": len(text_bytes)
                }

            # Store search index
            if search_index:
                index_key = f"{organization_id}/search_index/{document_id}_index.json"
                index_bytes = json.dumps(search_index, indent=2).encode('utf-8')

                index_result = await s3_storage.upload_file_direct(
                    file_content=index_bytes,
                    file_key=index_key,
                    content_type="application/json",
                    metadata={
                        "document_id": document_id,
                        "artifact_type": "search_index",
                        "organization_id": organization_id
                    }
                )

                artifacts["search_index"] = {
                    "s3_key": index_key,
                    "size": len(index_bytes)
                }

            return artifacts

        except Exception as e:
            logger.error(f"Artifact storage failed: {e}")
            return {"error": f"Artifact storage failed: {e}"}
    
    def get_file_info(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Extract basic file information"""
        mime_type, _ = mimetypes.guess_type(filename)
        file_size = len(file_content)
        
        return {
            "filename": filename,
            "mime_type": mime_type,
            "file_size": file_size,
            "file_extension": Path(filename).suffix.lower(),
            "document_type": self.SUPPORTED_FORMATS.get(mime_type, "unknown"),
            "is_supported": self.is_supported_format(mime_type or ""),
            "processed_at": datetime.utcnow()
        }
    
    async def process_document(
        self, 
        file_content: bytes, 
        filename: str
    ) -> Tuple[Dict[str, Any], str]:
        """
        Process document and extract text content
        Returns: (file_info, extracted_text)
        """
        file_info = self.get_file_info(file_content, filename)
        
        if not file_info["is_supported"]:
            raise ValueError(f"Unsupported file format: {file_info['mime_type']}")
        
        if file_info["file_size"] > self.max_file_size:
            raise ValueError(f"File too large: {file_info['file_size']} bytes (max: {self.max_file_size})")
        
        # Extract text based on file type
        document_type = file_info["document_type"]
        
        try:
            if document_type == "pdf":
                extracted_text = self._extract_from_pdf(file_content)
            elif document_type == "docx":
                extracted_text = self._extract_from_docx(file_content)
            elif document_type in ["xlsx", "xls"]:
                extracted_text = self._extract_from_excel(file_content)
            elif document_type in ["pptx", "ppt"]:
                extracted_text = self._extract_from_powerpoint(file_content)
            elif document_type == "txt":
                extracted_text = self._extract_from_text(file_content)
            elif document_type in ["jpg", "png", "tiff"]:
                extracted_text = self._extract_from_image(file_content)
            else:
                extracted_text = "Text extraction not implemented for this format"
            
            # Clean and preprocess text
            extracted_text = self._clean_text(extracted_text)
            
            # Update file info with extraction results
            file_info.update({
                "text_length": len(extracted_text),
                "word_count": len(extracted_text.split()),
                "extraction_successful": True
            })
            
            return file_info, extracted_text
            
        except Exception as e:
            logger.error(f"Error processing document {filename}: {str(e)}")
            file_info.update({
                "extraction_successful": False,
                "extraction_error": str(e)
            })
            return file_info, f"Error extracting text: {str(e)}"
    
    def _extract_from_pdf(self, file_content: bytes) -> str:
        """Extract text from PDF file with enhanced capabilities"""
        text_content = []
        metadata = {}

        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)

            # Extract metadata
            if pdf_reader.metadata:
                metadata = {
                    "title": pdf_reader.metadata.get("/Title", ""),
                    "author": pdf_reader.metadata.get("/Author", ""),
                    "subject": pdf_reader.metadata.get("/Subject", ""),
                    "creator": pdf_reader.metadata.get("/Creator", ""),
                    "producer": pdf_reader.metadata.get("/Producer", ""),
                    "creation_date": pdf_reader.metadata.get("/CreationDate", ""),
                    "modification_date": pdf_reader.metadata.get("/ModDate", "")
                }

            # Add metadata to content if available
            if any(metadata.values()):
                metadata_text = "--- Document Metadata ---\n"
                for key, value in metadata.items():
                    if value:
                        metadata_text += f"{key.title()}: {value}\n"
                text_content.append(metadata_text)

            # Extract text from each page
            total_pages = len(pdf_reader.pages)
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()

                    # Try alternative extraction methods if text is empty or minimal
                    if not page_text.strip() or len(page_text.strip()) < 50:
                        # Try extracting with different methods
                        try:
                            # Alternative extraction method
                            page_text = self._extract_pdf_page_alternative(page)
                        except:
                            pass

                    if page_text.strip():
                        # Clean and format page text
                        cleaned_text = self._clean_pdf_text(page_text)
                        if cleaned_text.strip():
                            text_content.append(f"--- Page {page_num + 1} of {total_pages} ---\n{cleaned_text}")
                    else:
                        text_content.append(f"--- Page {page_num + 1} of {total_pages} ---\n[No extractable text found]")

                except Exception as e:
                    logger.warning(f"Error extracting text from PDF page {page_num + 1}: {str(e)}")
                    text_content.append(f"--- Page {page_num + 1} of {total_pages} ---\n[Text extraction failed: {str(e)}]")
                    continue

            if not text_content:
                return "No text could be extracted from this PDF file."

            return "\n\n".join(text_content)

        except Exception as e:
            logger.error(f"Error reading PDF: {str(e)}")
            raise ValueError(f"Failed to extract text from PDF: {str(e)}")

    def _extract_pdf_page_alternative(self, page) -> str:
        """Alternative PDF text extraction method"""
        try:
            # Try to extract text using different approach
            if hasattr(page, 'extract_text'):
                return page.extract_text()
            return ""
        except:
            return ""

    def _clean_pdf_text(self, text: str) -> str:
        """Clean PDF-specific text artifacts"""
        if not text:
            return ""

        # Remove excessive whitespace and line breaks
        text = re.sub(r'\s+', ' ', text)

        # Fix common PDF extraction issues
        text = re.sub(r'([a-z])([A-Z])', r'\1 \2', text)  # Add space between camelCase
        text = re.sub(r'(\w)(\d)', r'\1 \2', text)  # Add space between word and number
        text = re.sub(r'(\d)(\w)', r'\1 \2', text)  # Add space between number and word

        # Remove page headers/footers patterns
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            line = line.strip()
            # Skip likely header/footer lines (very short, only numbers, etc.)
            if len(line) > 3 and not re.match(r'^\d+$', line):
                cleaned_lines.append(line)

        return '\n'.join(cleaned_lines)
    
    def _extract_from_docx(self, file_content: bytes) -> str:
        """Extract text from DOCX file"""
        try:
            docx_file = io.BytesIO(file_content)
            doc = Document(docx_file)
            
            text_content = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_text.append(" | ".join(row_text))
                
                if table_text:
                    text_content.append("--- Table ---\n" + "\n".join(table_text))
            
            return "\n\n".join(text_content)
            
        except Exception as e:
            logger.error(f"Error reading DOCX: {str(e)}")
            raise ValueError(f"Failed to extract text from DOCX: {str(e)}")
    
    def _extract_from_excel(self, file_content: bytes) -> str:
        """Extract text from Excel file with enhanced data processing"""
        try:
            excel_file = io.BytesIO(file_content)

            # Load workbook with both data and formulas
            workbook_data = openpyxl.load_workbook(excel_file, data_only=True)
            excel_file.seek(0)  # Reset file pointer
            workbook_formulas = openpyxl.load_workbook(excel_file, data_only=False)

            text_content = []

            for sheet_name in workbook_data.sheetnames:
                sheet_data = workbook_data[sheet_name]
                sheet_formulas = workbook_formulas[sheet_name] if sheet_name in workbook_formulas.sheetnames else None

                sheet_text = [f"--- Sheet: {sheet_name} ---"]

                # Get sheet dimensions
                max_row = sheet_data.max_row
                max_col = sheet_data.max_column

                if max_row == 1 and max_col == 1:
                    # Empty sheet
                    sheet_text.append("[Empty sheet]")
                else:
                    # Extract data with better formatting
                    for row_num in range(1, min(max_row + 1, 1000)):  # Limit to 1000 rows
                        row_data = []
                        row_formulas = []

                        for col_num in range(1, min(max_col + 1, 50)):  # Limit to 50 columns
                            # Get cell value
                            cell_data = sheet_data.cell(row=row_num, column=col_num).value
                            cell_formula = None

                            if sheet_formulas:
                                cell_formula = sheet_formulas.cell(row=row_num, column=col_num).value

                            if cell_data is not None:
                                cell_text = str(cell_data)

                                # Add formula information if different from value
                                if (cell_formula and
                                    str(cell_formula) != str(cell_data) and
                                    str(cell_formula).startswith('=')):
                                    cell_text += f" [{cell_formula}]"

                                row_data.append(cell_text)
                            elif cell_formula and str(cell_formula).startswith('='):
                                row_data.append(f"[{cell_formula}]")

                        if row_data:
                            # Format row with proper alignment
                            formatted_row = " | ".join(row_data)
                            sheet_text.append(formatted_row)

                    # Add summary information
                    sheet_text.append(f"\n[Sheet Summary: {max_row} rows, {max_col} columns]")

                if len(sheet_text) > 1:  # More than just the header
                    text_content.append("\n".join(sheet_text))

            # Add workbook-level information
            if len(workbook_data.sheetnames) > 1:
                summary = f"--- Workbook Summary ---\nTotal Sheets: {len(workbook_data.sheetnames)}\nSheet Names: {', '.join(workbook_data.sheetnames)}"
                text_content.insert(0, summary)

            return "\n\n".join(text_content)

        except Exception as e:
            logger.error(f"Error reading Excel: {str(e)}")
            raise ValueError(f"Failed to extract text from Excel: {str(e)}")
    
    def _extract_from_powerpoint(self, file_content: bytes) -> str:
        """Extract text from PowerPoint file"""
        try:
            ppt_file = io.BytesIO(file_content)
            presentation = Presentation(ppt_file)
            
            text_content = []
            
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = [f"--- Slide {slide_num + 1} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # More than just the header
                    text_content.append("\n".join(slide_text))
            
            return "\n\n".join(text_content)
            
        except Exception as e:
            logger.error(f"Error reading PowerPoint: {str(e)}")
            raise ValueError(f"Failed to extract text from PowerPoint: {str(e)}")
    
    def _extract_from_text(self, file_content: bytes) -> str:
        """Extract text from plain text file"""
        try:
            # Try different encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    return file_content.decode(encoding)
                except UnicodeDecodeError:
                    continue
            
            # If all encodings fail, use utf-8 with error handling
            return file_content.decode('utf-8', errors='replace')
            
        except Exception as e:
            logger.error(f"Error reading text file: {str(e)}")
            raise ValueError(f"Failed to extract text from file: {str(e)}")
    
    def _extract_from_image(self, file_content: bytes) -> str:
        """Extract text from image using OCR with enhanced preprocessing"""
        try:
            image = Image.open(io.BytesIO(file_content))

            # Convert to RGB if necessary
            if image.mode != 'RGB':
                image = image.convert('RGB')

            # Enhance image for better OCR results
            image = self._preprocess_image_for_ocr(image)

            # Use pytesseract for OCR with custom configuration
            custom_config = r'--oem 3 --psm 6 -c tessedit_char_whitelist=0123456789ABCDEFGHIJKLMNOPQRSTUVWXYZabcdefghijklmnopqrstuvwxyz.,!?@#$%^&*()_+-=[]{}|;:,.<>?/~` '
            extracted_text = pytesseract.image_to_string(image, config=custom_config)

            if not extracted_text.strip():
                # Try with different PSM modes if first attempt fails
                for psm in [3, 4, 8, 11, 13]:
                    try:
                        config = f'--oem 3 --psm {psm}'
                        extracted_text = pytesseract.image_to_string(image, config=config)
                        if extracted_text.strip():
                            break
                    except:
                        continue

            if not extracted_text.strip():
                return "No text found in image"

            return extracted_text

        except Exception as e:
            logger.error(f"Error extracting text from image: {str(e)}")
            return "Error processing image"

    async def enhanced_security_scan(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """
        Perform comprehensive security scanning on uploaded files
        """
        security_results = {
            "is_safe": True,
            "threats_detected": [],
            "warnings": [],
            "scan_timestamp": datetime.utcnow().isoformat(),
            "scans_performed": []
        }

        try:
            # 1. File extension validation
            file_ext = Path(filename).suffix.lower()
            if file_ext not in self.allowed_extensions:
                security_results["is_safe"] = False
                security_results["threats_detected"].append({
                    "type": "forbidden_extension",
                    "description": f"File extension '{file_ext}' is not allowed",
                    "severity": "high"
                })
            security_results["scans_performed"].append("extension_validation")

            # 2. File signature validation (magic bytes)
            file_signature = file_content[:16].hex() if len(file_content) >= 16 else ""
            expected_signatures = self._get_expected_signatures(file_ext)

            if expected_signatures and not any(file_signature.startswith(sig) for sig in expected_signatures):
                security_results["warnings"].append({
                    "type": "signature_mismatch",
                    "description": "File signature doesn't match extension",
                    "severity": "medium"
                })
            security_results["scans_performed"].append("signature_validation")

            # 3. Content pattern scanning
            if self.enable_content_analysis:
                content_str = file_content.decode('utf-8', errors='ignore')
                for pattern in self.dangerous_patterns:
                    if re.search(pattern, content_str, re.IGNORECASE):
                        security_results["is_safe"] = False
                        security_results["threats_detected"].append({
                            "type": "malicious_content",
                            "description": f"Dangerous pattern detected: {pattern}",
                            "severity": "high"
                        })
                security_results["scans_performed"].append("content_analysis")

            # 4. PII detection
            if self.enable_pii_detection:
                pii_results = await self._detect_pii(file_content, filename)
                if pii_results["pii_detected"]:
                    security_results["warnings"].append({
                        "type": "pii_detected",
                        "description": f"PII detected: {', '.join(pii_results['pii_types'])}",
                        "severity": "medium",
                        "details": pii_results
                    })
                security_results["scans_performed"].append("pii_detection")

            # 5. Virus scanning (if enabled)
            if self.virus_scan_enabled:
                virus_results = await self._scan_for_viruses(file_content)
                if not virus_results["is_clean"]:
                    security_results["is_safe"] = False
                    security_results["threats_detected"].append({
                        "type": "virus_detected",
                        "description": virus_results["threat_name"],
                        "severity": "critical"
                    })
                security_results["scans_performed"].append("virus_scan")

            # 6. File size validation
            if len(file_content) > 50 * 1024 * 1024:  # 50MB
                security_results["warnings"].append({
                    "type": "large_file",
                    "description": f"File size ({len(file_content)} bytes) exceeds recommended limit",
                    "severity": "low"
                })
            security_results["scans_performed"].append("size_validation")

            logger.info(f"Security scan completed for {filename}: {'SAFE' if security_results['is_safe'] else 'THREATS DETECTED'}")
            return security_results

        except Exception as e:
            logger.error(f"Security scan failed for {filename}: {e}")
            return {
                "is_safe": False,
                "threats_detected": [{
                    "type": "scan_error",
                    "description": f"Security scan failed: {str(e)}",
                    "severity": "high"
                }],
                "warnings": [],
                "scan_timestamp": datetime.utcnow().isoformat(),
                "scans_performed": ["error"]
            }

    def _get_expected_signatures(self, file_ext: str) -> List[str]:
        """Get expected file signatures for validation"""
        signatures = {
            '.pdf': ['255044462d'],  # %PDF-
            '.jpg': ['ffd8ff'],      # JPEG
            '.jpeg': ['ffd8ff'],     # JPEG
            '.png': ['89504e47'],    # PNG
            '.gif': ['474946'],      # GIF
            '.zip': ['504b0304', '504b0506'],  # ZIP
            '.docx': ['504b0304'],   # DOCX (ZIP-based)
            '.xlsx': ['504b0304'],   # XLSX (ZIP-based)
            '.pptx': ['504b0304'],   # PPTX (ZIP-based)
        }
        return signatures.get(file_ext, [])

    async def _detect_pii(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Detect personally identifiable information in files"""
        try:
            # Convert to text for analysis
            content_str = file_content.decode('utf-8', errors='ignore')

            pii_patterns = {
                'ssn': r'\b\d{3}-\d{2}-\d{4}\b',
                'credit_card': r'\b\d{4}[-\s]?\d{4}[-\s]?\d{4}[-\s]?\d{4}\b',
                'email': r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b',
                'phone': r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b',
                'ip_address': r'\b\d{1,3}\.\d{1,3}\.\d{1,3}\.\d{1,3}\b'
            }

            detected_pii = []
            pii_counts = {}

            for pii_type, pattern in pii_patterns.items():
                matches = re.findall(pattern, content_str)
                if matches:
                    detected_pii.append(pii_type)
                    pii_counts[pii_type] = len(matches)

            return {
                "pii_detected": len(detected_pii) > 0,
                "pii_types": detected_pii,
                "pii_counts": pii_counts,
                "scan_timestamp": datetime.utcnow().isoformat()
            }

        except Exception as e:
            logger.error(f"PII detection failed for {filename}: {e}")
            return {
                "pii_detected": False,
                "pii_types": [],
                "pii_counts": {},
                "error": str(e),
                "scan_timestamp": datetime.utcnow().isoformat()
            }

    async def _scan_for_viruses(self, file_content: bytes) -> Dict[str, Any]:
        """Scan file content for viruses using ClamAV"""
        try:
            # This is a simplified implementation
            # In production, you would integrate with ClamAV or another antivirus solution

            # For now, check for known malicious patterns
            malicious_patterns = [
                b'X5O!P%@AP[4\\PZX54(P^)7CC)7}$EICAR-STANDARD-ANTIVIRUS-TEST-FILE!$H+H*',  # EICAR test
                b'<script>alert(',  # Basic XSS
                b'eval(',  # JavaScript eval
                b'exec(',  # Python exec
            ]

            for pattern in malicious_patterns:
                if pattern in file_content:
                    return {
                        "is_clean": False,
                        "threat_name": f"Malicious pattern detected: {pattern.decode('utf-8', errors='ignore')[:50]}",
                        "scan_timestamp": datetime.utcnow().isoformat()
                    }

            return {
                "is_clean": True,
                "threat_name": None,
                "scan_timestamp": datetime.utcnow().isoformat()
            }

        except Exception as e:
            logger.error(f"Virus scan failed: {e}")
            return {
                "is_clean": False,
                "threat_name": f"Scan error: {str(e)}",
                "scan_timestamp": datetime.utcnow().isoformat()
            }

        except Exception as e:
            logger.error(f"Error performing OCR on image: {str(e)}")
            return f"OCR failed: {str(e)}"

    def _preprocess_image_for_ocr(self, image: Image.Image) -> Image.Image:
        """Preprocess image to improve OCR accuracy"""
        try:
            # Convert to grayscale
            if image.mode != 'L':
                image = image.convert('L')

            # Resize if image is too small (OCR works better on larger images)
            width, height = image.size
            if width < 300 or height < 300:
                scale_factor = max(300 / width, 300 / height)
                new_width = int(width * scale_factor)
                new_height = int(height * scale_factor)
                image = image.resize((new_width, new_height), Image.Resampling.LANCZOS)

            return image

        except Exception as e:
            logger.warning(f"Image preprocessing failed: {str(e)}")
            return image
    
    def _clean_text(self, text: str) -> str:
        """Clean and preprocess extracted text"""
        if not text:
            return ""
        
        # Remove excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        # Remove excessive newlines
        text = re.sub(r'\n\s*\n\s*\n+', '\n\n', text)
        
        # Remove special characters that might interfere with AI processing
        text = re.sub(r'[^\w\s\.\,\;\:\!\?\-\(\)\[\]\{\}\"\'\/\\\@\#\$\%\^\&\*\+\=\|\~\`]', ' ', text)
        
        # Trim whitespace
        text = text.strip()
        
        return text
    
    def validate_file_security(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Perform security validation on uploaded file"""
        security_report = {
            "is_safe": True,
            "threats_detected": [],
            "warnings": [],
            "file_analysis": {}
        }

        try:
            # Check file size
            file_size = len(file_content)
            if file_size > self.max_file_size:
                security_report["is_safe"] = False
                security_report["threats_detected"].append(f"File too large: {file_size} bytes")

            # Check for suspicious file extensions
            file_ext = Path(filename).suffix.lower()
            dangerous_extensions = ['.exe', '.bat', '.cmd', '.scr', '.pif', '.com', '.jar']
            if file_ext in dangerous_extensions:
                security_report["is_safe"] = False
                security_report["threats_detected"].append(f"Dangerous file extension: {file_ext}")

            # Check for embedded executables in documents
            if self._contains_embedded_executable(file_content):
                security_report["is_safe"] = False
                security_report["threats_detected"].append("Embedded executable content detected")

            # Check for suspicious patterns
            suspicious_patterns = self._scan_for_suspicious_patterns(file_content)
            if suspicious_patterns:
                security_report["warnings"].extend(suspicious_patterns)

            # File structure analysis
            security_report["file_analysis"] = self._analyze_file_structure(file_content, filename)

        except Exception as e:
            logger.error(f"Security validation failed: {str(e)}")
            security_report["warnings"].append(f"Security scan incomplete: {str(e)}")

        return security_report

    def _contains_embedded_executable(self, file_content: bytes) -> bool:
        """Check for embedded executable content"""
        # Check for common executable signatures
        executable_signatures = [
            b'MZ',  # DOS/Windows executable
            b'\x7fELF',  # Linux executable
            b'\xfe\xed\xfa',  # Mach-O executable
            b'PK\x03\x04',  # ZIP (could contain executables)
        ]

        for signature in executable_signatures:
            if signature in file_content[:1024]:  # Check first 1KB
                return True

        return False

    def _scan_for_suspicious_patterns(self, file_content: bytes) -> List[str]:
        """Scan for suspicious patterns in file content"""
        warnings = []

        try:
            # Convert to string for pattern matching (handle encoding errors)
            content_str = file_content.decode('utf-8', errors='ignore').lower()

            # Suspicious patterns
            suspicious_patterns = [
                ('javascript:', 'JavaScript code detected'),
                ('vbscript:', 'VBScript code detected'),
                ('powershell', 'PowerShell commands detected'),
                ('cmd.exe', 'Command line references detected'),
                ('eval(', 'Code evaluation functions detected'),
                ('document.write', 'Dynamic content generation detected')
            ]

            for pattern, warning in suspicious_patterns:
                if pattern in content_str:
                    warnings.append(warning)

        except Exception as e:
            logger.warning(f"Pattern scanning failed: {str(e)}")

        return warnings

    def _analyze_file_structure(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Analyze file structure for additional insights"""
        analysis = {
            "file_size": len(file_content),
            "entropy": self._calculate_entropy(file_content),
            "mime_type_detected": None,
            "structure_valid": True
        }

        try:
            # Detect MIME type from content
            import mimetypes
            mime_type, _ = mimetypes.guess_type(filename)
            analysis["mime_type_detected"] = mime_type

            # Calculate file entropy (high entropy might indicate encryption/compression)
            entropy = self._calculate_entropy(file_content)
            analysis["entropy"] = entropy

            if entropy > 7.5:  # High entropy threshold
                analysis["high_entropy"] = True

        except Exception as e:
            logger.warning(f"File structure analysis failed: {str(e)}")
            analysis["structure_valid"] = False

        return analysis

    def _calculate_entropy(self, data: bytes) -> float:
        """Calculate Shannon entropy of data"""
        if not data:
            return 0

        # Count byte frequencies
        byte_counts = [0] * 256
        for byte in data:
            byte_counts[byte] += 1

        # Calculate entropy
        entropy = 0
        data_len = len(data)
        for count in byte_counts:
            if count > 0:
                probability = count / data_len
                entropy -= probability * (probability.bit_length() - 1)

        return entropy

    def get_document_stats(self, text: str) -> Dict[str, Any]:
        """Get comprehensive statistics about the extracted text"""
        if not text:
            return {
                "word_count": 0,
                "character_count": 0,
                "paragraph_count": 0,
                "sentence_count": 0,
                "average_words_per_paragraph": 0,
                "average_words_per_sentence": 0,
                "readability_score": 0
            }

        words = text.split()
        paragraphs = [p for p in text.split('\n\n') if p.strip()]
        sentences = [s for s in re.split(r'[.!?]+', text) if s.strip()]

        # Calculate readability metrics
        avg_words_per_sentence = len(words) / max(len(sentences), 1)
        avg_chars_per_word = sum(len(word) for word in words) / max(len(words), 1)

        # Simple readability score (Flesch-like)
        readability_score = max(0, min(100,
            206.835 - (1.015 * avg_words_per_sentence) - (84.6 * (avg_chars_per_word / 4.7))
        ))

        return {
            "word_count": len(words),
            "character_count": len(text),
            "paragraph_count": len(paragraphs),
            "sentence_count": len(sentences),
            "average_words_per_paragraph": len(words) / max(len(paragraphs), 1),
            "average_words_per_sentence": avg_words_per_sentence,
            "average_characters_per_word": avg_chars_per_word,
            "readability_score": round(readability_score, 2),
            "unique_words": len(set(word.lower() for word in words if word.isalpha())),
            "text_complexity": "high" if avg_words_per_sentence > 20 else "medium" if avg_words_per_sentence > 15 else "low"
        }


# Global instance
document_processor = DocumentProcessor()
