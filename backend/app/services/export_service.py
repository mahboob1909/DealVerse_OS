#!/usr/bin/env python3
"""
Export Service for DealVerse OS
Handles PDF, Excel, and PowerPoint export functionality
"""
import io
import json
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any, Optional, Union
from uuid import UUID

import pandas as pd
from fastapi import HTTPException
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
from reportlab.platypus.flowables import PageBreak
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import LineChart, BarChart, Reference
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from app.core.config import settings
import logging

logger = logging.getLogger(__name__)


class ExportService:
    """Service for handling various export formats"""
    
    def __init__(self):
        self.styles = getSampleStyleSheet()
        self.brand_colors = {
            'primary': '#1a2332',
            'secondary': '#0066ff', 
            'accent': '#00c896',
            'text': '#333333',
            'light_gray': '#f8f9fa'
        }
        
    async def export_financial_model_to_excel(
        self, 
        model_data: Dict[str, Any],
        model_id: UUID,
        organization_name: str = "DealVerse Organization"
    ) -> bytes:
        """Export financial model to Excel format"""
        try:
            wb = Workbook()
            
            # Remove default sheet and create custom sheets
            wb.remove(wb.active)
            
            # Create Summary sheet
            summary_ws = wb.create_sheet("Executive Summary")
            self._create_excel_summary_sheet(summary_ws, model_data, organization_name)
            
            # Create Financial Projections sheet
            projections_ws = wb.create_sheet("Financial Projections")
            self._create_excel_projections_sheet(projections_ws, model_data)
            
            # Create Assumptions sheet
            assumptions_ws = wb.create_sheet("Key Assumptions")
            self._create_excel_assumptions_sheet(assumptions_ws, model_data)
            
            # Create Charts sheet
            charts_ws = wb.create_sheet("Charts & Analysis")
            self._create_excel_charts_sheet(charts_ws, model_data)
            
            # Save to bytes
            excel_buffer = io.BytesIO()
            wb.save(excel_buffer)
            excel_buffer.seek(0)
            
            return excel_buffer.getvalue()
            
        except Exception as e:
            logger.error(f"Excel export failed for model {model_id}: {e}")
            raise HTTPException(status_code=500, detail=f"Excel export failed: {str(e)}")
    
    def _create_excel_summary_sheet(self, ws, model_data: Dict, organization_name: str):
        """Create executive summary sheet in Excel"""
        # Header styling
        header_font = Font(name='Arial', size=16, bold=True, color='FFFFFF')
        header_fill = PatternFill(start_color='1a2332', end_color='1a2332', fill_type='solid')
        
        # Title
        ws['A1'] = f"{organization_name} - Financial Model Summary"
        ws['A1'].font = header_font
        ws['A1'].fill = header_fill
        ws.merge_cells('A1:F1')
        
        # Model details
        ws['A3'] = "Model Name:"
        ws['B3'] = model_data.get('name', 'Untitled Model')
        ws['A4'] = "Created Date:"
        ws['B4'] = model_data.get('created_at', datetime.now().strftime('%Y-%m-%d'))
        ws['A5'] = "Last Updated:"
        ws['B5'] = model_data.get('updated_at', datetime.now().strftime('%Y-%m-%d'))
        
        # Key metrics
        ws['A7'] = "Key Financial Metrics"
        ws['A7'].font = Font(bold=True, size=14)
        
        metrics = model_data.get('key_metrics', {})
        row = 8
        for metric, value in metrics.items():
            ws[f'A{row}'] = metric.replace('_', ' ').title()
            ws[f'B{row}'] = value
            row += 1
    
    def _create_excel_projections_sheet(self, ws, model_data: Dict):
        """Create financial projections sheet"""
        projections = model_data.get('projections', {})
        
        # Headers
        headers = ['Year', 'Revenue', 'COGS', 'Gross Profit', 'Operating Expenses', 'EBITDA', 'Net Income']
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = Font(bold=True)
            cell.fill = PatternFill(start_color='0066ff', end_color='0066ff', fill_type='solid')
            cell.font = Font(color='FFFFFF', bold=True)
        
        # Data rows
        years = projections.get('years', [])
        for row, year_data in enumerate(years, 2):
            ws.cell(row=row, column=1, value=year_data.get('year', ''))
            ws.cell(row=row, column=2, value=year_data.get('revenue', 0))
            ws.cell(row=row, column=3, value=year_data.get('cogs', 0))
            ws.cell(row=row, column=4, value=year_data.get('gross_profit', 0))
            ws.cell(row=row, column=5, value=year_data.get('operating_expenses', 0))
            ws.cell(row=row, column=6, value=year_data.get('ebitda', 0))
            ws.cell(row=row, column=7, value=year_data.get('net_income', 0))
    
    def _create_excel_assumptions_sheet(self, ws, model_data: Dict):
        """Create assumptions sheet"""
        assumptions = model_data.get('assumptions', {})
        
        ws['A1'] = "Key Model Assumptions"
        ws['A1'].font = Font(bold=True, size=14)
        
        row = 3
        for category, items in assumptions.items():
            ws[f'A{row}'] = category.replace('_', ' ').title()
            ws[f'A{row}'].font = Font(bold=True)
            row += 1
            
            if isinstance(items, dict):
                for key, value in items.items():
                    ws[f'B{row}'] = key.replace('_', ' ').title()
                    ws[f'C{row}'] = value
                    row += 1
            row += 1
    
    def _create_excel_charts_sheet(self, ws, model_data: Dict):
        """Create charts and analysis sheet"""
        ws['A1'] = "Financial Analysis & Charts"
        ws['A1'].font = Font(bold=True, size=14)
        
        # Add chart data
        projections = model_data.get('projections', {})
        years = projections.get('years', [])
        
        if years:
            # Revenue chart data
            ws['A3'] = "Year"
            ws['B3'] = "Revenue"
            ws['C3'] = "Net Income"
            
            for row, year_data in enumerate(years, 4):
                ws[f'A{row}'] = year_data.get('year', '')
                ws[f'B{row}'] = year_data.get('revenue', 0)
                ws[f'C{row}'] = year_data.get('net_income', 0)
            
            # Create line chart
            chart = LineChart()
            chart.title = "Revenue & Net Income Projection"
            chart.style = 10
            chart.x_axis.title = 'Year'
            chart.y_axis.title = 'Amount ($)'
            
            data = Reference(ws, min_col=2, min_row=3, max_col=3, max_row=3+len(years))
            cats = Reference(ws, min_col=1, min_row=4, max_row=3+len(years))
            chart.add_data(data, titles_from_data=True)
            chart.set_categories(cats)
            
            ws.add_chart(chart, "E3")
    
    async def export_financial_model_to_pdf(
        self,
        model_data: Dict[str, Any],
        model_id: UUID,
        organization_name: str = "DealVerse Organization"
    ) -> bytes:
        """Export financial model to PDF format"""
        try:
            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=A4)
            story = []
            
            # Title
            title_style = ParagraphStyle(
                'CustomTitle',
                parent=self.styles['Heading1'],
                fontSize=24,
                spaceAfter=30,
                textColor=colors.HexColor(self.brand_colors['primary'])
            )
            
            story.append(Paragraph(f"{organization_name}<br/>Financial Model Report", title_style))
            story.append(Spacer(1, 20))
            
            # Model information
            model_info = [
                ['Model Name:', model_data.get('name', 'Untitled Model')],
                ['Created Date:', model_data.get('created_at', datetime.now().strftime('%Y-%m-%d'))],
                ['Last Updated:', model_data.get('updated_at', datetime.now().strftime('%Y-%m-%d'))],
                ['Model Type:', model_data.get('model_type', 'Financial Projection')]
            ]
            
            info_table = Table(model_info, colWidths=[2*inch, 4*inch])
            info_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (0, -1), colors.HexColor(self.brand_colors['light_gray'])),
                ('TEXTCOLOR', (0, 0), (-1, -1), colors.black),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 10),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))
            
            story.append(info_table)
            story.append(Spacer(1, 30))
            
            # Key metrics section
            story.append(Paragraph("Key Financial Metrics", self.styles['Heading2']))
            story.append(Spacer(1, 12))
            
            metrics = model_data.get('key_metrics', {})
            if metrics:
                metrics_data = [['Metric', 'Value']]
                for metric, value in metrics.items():
                    metrics_data.append([metric.replace('_', ' ').title(), str(value)])
                
                metrics_table = Table(metrics_data, colWidths=[3*inch, 2*inch])
                metrics_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(self.brand_colors['secondary'])),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, -1), 10),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                
                story.append(metrics_table)
            
            story.append(PageBreak())
            
            # Financial projections
            story.append(Paragraph("Financial Projections", self.styles['Heading2']))
            story.append(Spacer(1, 12))
            
            projections = model_data.get('projections', {})
            years = projections.get('years', [])
            
            if years:
                proj_data = [['Year', 'Revenue', 'COGS', 'Gross Profit', 'Operating Expenses', 'EBITDA', 'Net Income']]
                for year_data in years:
                    proj_data.append([
                        str(year_data.get('year', '')),
                        f"${year_data.get('revenue', 0):,.0f}",
                        f"${year_data.get('cogs', 0):,.0f}",
                        f"${year_data.get('gross_profit', 0):,.0f}",
                        f"${year_data.get('operating_expenses', 0):,.0f}",
                        f"${year_data.get('ebitda', 0):,.0f}",
                        f"${year_data.get('net_income', 0):,.0f}"
                    ])
                
                proj_table = Table(proj_data, colWidths=[0.8*inch, 1*inch, 1*inch, 1*inch, 1.2*inch, 1*inch, 1*inch])
                proj_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(self.brand_colors['primary'])),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, -1), 8),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor(self.brand_colors['light_gray'])])
                ]))
                
                story.append(proj_table)
            
            # Build PDF
            doc.build(story)
            buffer.seek(0)
            
            return buffer.getvalue()
            
        except Exception as e:
            logger.error(f"PDF export failed for model {model_id}: {e}")
            raise HTTPException(status_code=500, detail=f"PDF export failed: {str(e)}")


    async def export_presentation_to_pptx(
        self,
        presentation_data: Dict[str, Any],
        presentation_id: UUID,
        organization_name: str = "DealVerse Organization"
    ) -> bytes:
        """Export presentation to PowerPoint format"""
        try:
            prs = Presentation()

            # Set slide size to widescreen
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)

            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = presentation_data.get('title', 'DealVerse Presentation')
            subtitle.text = f"{organization_name}\n{datetime.now().strftime('%B %Y')}"

            # Style title slide
            title.text_frame.paragraphs[0].font.size = Pt(44)
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(26, 35, 50)  # brand primary

            # Add slides from presentation data
            slides_data = presentation_data.get('slides', [])
            for slide_data in slides_data:
                self._create_pptx_slide(prs, slide_data)

            # Save to bytes
            pptx_buffer = io.BytesIO()
            prs.save(pptx_buffer)
            pptx_buffer.seek(0)

            return pptx_buffer.getvalue()

        except Exception as e:
            logger.error(f"PowerPoint export failed for presentation {presentation_id}: {e}")
            raise HTTPException(status_code=500, detail=f"PowerPoint export failed: {str(e)}")

    def _create_pptx_slide(self, prs: Presentation, slide_data: Dict[str, Any]):
        """Create a PowerPoint slide from slide data"""
        slide_type = slide_data.get('type', 'content')

        if slide_type == 'title_content':
            layout = prs.slide_layouts[1]  # Title and Content
        elif slide_type == 'two_content':
            layout = prs.slide_layouts[3]  # Two Content
        elif slide_type == 'blank':
            layout = prs.slide_layouts[6]  # Blank
        else:
            layout = prs.slide_layouts[1]  # Default to Title and Content

        slide = prs.slides.add_slide(layout)

        # Add title
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = slide_data.get('title', '')
            slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
            slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(26, 35, 50)

        # Add content
        content = slide_data.get('content', '')
        if content and len(slide.placeholders) > 1:
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = content

            # Style content
            for paragraph in content_placeholder.text_frame.paragraphs:
                paragraph.font.size = Pt(18)
                paragraph.font.color.rgb = RGBColor(51, 51, 51)

        # Add charts or tables if present
        charts = slide_data.get('charts', [])
        for chart_data in charts:
            self._add_chart_to_slide(slide, chart_data)

    def _add_chart_to_slide(self, slide, chart_data: Dict[str, Any]):
        """Add a chart to a PowerPoint slide"""
        # This is a simplified chart implementation
        # In a full implementation, you would use python-pptx chart capabilities
        chart_title = chart_data.get('title', 'Chart')
        chart_type = chart_data.get('type', 'line')

        # Add a text box with chart information for now
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(2)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = f"Chart: {chart_title}\nType: {chart_type}\nData visualization would appear here"

    async def export_compliance_report_to_pdf(
        self,
        compliance_data: Dict[str, Any],
        organization_name: str = "DealVerse Organization"
    ) -> bytes:
        """Export compliance report to PDF format"""
        try:
            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=A4)
            story = []

            # Title
            title_style = ParagraphStyle(
                'ComplianceTitle',
                parent=self.styles['Heading1'],
                fontSize=24,
                spaceAfter=30,
                textColor=colors.HexColor(self.brand_colors['primary'])
            )

            story.append(Paragraph(f"{organization_name}<br/>Compliance Audit Report", title_style))
            story.append(Spacer(1, 20))

            # Report metadata
            report_info = [
                ['Report Date:', datetime.now().strftime('%Y-%m-%d %H:%M:%S')],
                ['Audit Period:', compliance_data.get('audit_period', 'Last 30 days')],
                ['Compliance Areas:', ', '.join(compliance_data.get('compliance_areas', ['All']))],
                ['Report Type:', compliance_data.get('report_type', 'Comprehensive')]
            ]

            info_table = Table(report_info, colWidths=[2*inch, 4*inch])
            info_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (0, -1), colors.HexColor(self.brand_colors['light_gray'])),
                ('TEXTCOLOR', (0, 0), (-1, -1), colors.black),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (0, -1), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 10),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))

            story.append(info_table)
            story.append(Spacer(1, 30))

            # Compliance summary
            story.append(Paragraph("Compliance Summary", self.styles['Heading2']))
            story.append(Spacer(1, 12))

            summary = compliance_data.get('summary', {})
            summary_data = [
                ['Total Checks Performed', str(summary.get('total_checks', 0))],
                ['Passed', str(summary.get('passed', 0))],
                ['Failed', str(summary.get('failed', 0))],
                ['Warnings', str(summary.get('warnings', 0))],
                ['Compliance Score', f"{summary.get('compliance_score', 0):.1f}%"]
            ]

            summary_table = Table(summary_data, colWidths=[3*inch, 2*inch])
            summary_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(self.brand_colors['secondary'])),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 10),
                ('GRID', (0, 0), (-1, -1), 1, colors.black)
            ]))

            story.append(summary_table)
            story.append(PageBreak())

            # Detailed findings
            story.append(Paragraph("Detailed Findings", self.styles['Heading2']))
            story.append(Spacer(1, 12))

            findings = compliance_data.get('findings', [])
            if findings:
                findings_data = [['Category', 'Status', 'Description', 'Risk Level']]
                for finding in findings:
                    findings_data.append([
                        finding.get('category', ''),
                        finding.get('status', ''),
                        finding.get('description', ''),
                        finding.get('risk_level', '')
                    ])

                findings_table = Table(findings_data, colWidths=[1.5*inch, 1*inch, 3*inch, 1*inch])
                findings_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor(self.brand_colors['primary'])),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, -1), 8),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor(self.brand_colors['light_gray'])])
                ]))

                story.append(findings_table)

            # Build PDF
            doc.build(story)
            buffer.seek(0)

            return buffer.getvalue()

        except Exception as e:
            logger.error(f"Compliance PDF export failed: {e}")
            raise HTTPException(status_code=500, detail=f"Compliance PDF export failed: {str(e)}")

    async def export_analytics_to_excel(
        self,
        analytics_data: Dict[str, Any],
        organization_name: str = "DealVerse Organization"
    ) -> bytes:
        """Export analytics dashboard data to Excel format"""
        try:
            wb = Workbook()
            wb.remove(wb.active)

            # Dashboard Overview sheet
            overview_ws = wb.create_sheet("Dashboard Overview")
            self._create_analytics_overview_sheet(overview_ws, analytics_data, organization_name)

            # Performance Metrics sheet
            metrics_ws = wb.create_sheet("Performance Metrics")
            self._create_analytics_metrics_sheet(metrics_ws, analytics_data)

            # Trends Analysis sheet
            trends_ws = wb.create_sheet("Trends Analysis")
            self._create_analytics_trends_sheet(trends_ws, analytics_data)

            # Save to bytes
            excel_buffer = io.BytesIO()
            wb.save(excel_buffer)
            excel_buffer.seek(0)

            return excel_buffer.getvalue()

        except Exception as e:
            logger.error(f"Analytics Excel export failed: {e}")
            raise HTTPException(status_code=500, detail=f"Analytics Excel export failed: {str(e)}")

    def _create_analytics_overview_sheet(self, ws, analytics_data: Dict, organization_name: str):
        """Create analytics overview sheet"""
        # Header
        ws['A1'] = f"{organization_name} - Analytics Dashboard"
        ws['A1'].font = Font(name='Arial', size=16, bold=True)
        ws.merge_cells('A1:F1')

        # Key metrics
        ws['A3'] = "Key Performance Indicators"
        ws['A3'].font = Font(bold=True, size=14)

        kpis = analytics_data.get('kpis', {})
        row = 4
        for kpi, value in kpis.items():
            ws[f'A{row}'] = kpi.replace('_', ' ').title()
            ws[f'B{row}'] = value
            row += 1

    def _create_analytics_metrics_sheet(self, ws, analytics_data: Dict):
        """Create performance metrics sheet"""
        ws['A1'] = "Performance Metrics"
        ws['A1'].font = Font(bold=True, size=14)

        metrics = analytics_data.get('metrics', [])
        if metrics:
            headers = ['Metric', 'Current Value', 'Previous Value', 'Change %']
            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=3, column=col, value=header)
                cell.font = Font(bold=True)

            for row, metric in enumerate(metrics, 4):
                ws.cell(row=row, column=1, value=metric.get('name', ''))
                ws.cell(row=row, column=2, value=metric.get('current_value', 0))
                ws.cell(row=row, column=3, value=metric.get('previous_value', 0))
                ws.cell(row=row, column=4, value=metric.get('change_percent', 0))

    def _create_analytics_trends_sheet(self, ws, analytics_data: Dict):
        """Create trends analysis sheet"""
        ws['A1'] = "Trends Analysis"
        ws['A1'].font = Font(bold=True, size=14)

        trends = analytics_data.get('trends', [])
        if trends:
            headers = ['Date', 'Metric', 'Value']
            for col, header in enumerate(headers, 1):
                cell = ws.cell(row=3, column=col, value=header)
                cell.font = Font(bold=True)

            for row, trend in enumerate(trends, 4):
                ws.cell(row=row, column=1, value=trend.get('date', ''))
                ws.cell(row=row, column=2, value=trend.get('metric', ''))
                ws.cell(row=row, column=3, value=trend.get('value', 0))


# Create global instance
export_service = ExportService()
